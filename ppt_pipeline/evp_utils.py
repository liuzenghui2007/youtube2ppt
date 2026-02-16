"""evp 相关：解析帧时间、从原视频抽帧、多图合成 PDF。"""
from __future__ import annotations

import re
import subprocess
from pathlib import Path
from typing import Callable


EVP_TMP_DIR = ".extract-video-ppt-tmp-data"


def parse_evp_frame_timestamps(output_dir: Path) -> list[float]:
    """从 evp 临时目录中的帧文件名解析时间戳（秒）。"""
    tmp = Path(output_dir) / EVP_TMP_DIR
    if not tmp.is_dir():
        return []
    # 例: frame00:00:01-0.56.jpg -> 1.0 秒
    pattern = re.compile(r"frame(\d{2}):(\d{2}):(\d{2})[-.]")
    times = []
    for f in tmp.iterdir():
        if f.suffix.lower() not in (".jpg", ".jpeg", ".png"):
            continue
        m = pattern.match(f.stem)
        if m:
            h, m_, s = int(m.group(1)), int(m.group(2)), int(m.group(3))
            times.append(h * 3600 + m_ * 60 + s)
    return sorted(times)


def extract_frames_at_times(
    video_path: Path,
    times_sec: list[float],
    output_dir: Path,
    *,
    progress_callback: Callable[[int, int], None] | None = None,
) -> list[Path]:
    """在给定时间点从视频抽帧，保存为 frame_001.png 等，返回路径列表。progress_callback(current_1based, total)。"""
    output_dir = Path(output_dir).resolve()
    output_dir.mkdir(parents=True, exist_ok=True)
    total = len(times_sec)
    paths = []
    for i, t in enumerate(times_sec):
        out = output_dir / f"frame_{i + 1:03d}.png"
        r = subprocess.run(
            ["ffmpeg", "-y", "-ss", str(t), "-i", str(video_path), "-vframes", "1", "-q:v", "2", str(out)],
            capture_output=True,
        )
        if r.returncode == 0 and out.is_file():
            paths.append(out)
        if progress_callback and total > 0:
            progress_callback(i + 1, total)
    return paths


def frames_to_pdf(image_paths: list[Path], pdf_path: Path, *, size_wh: tuple[int, int] | None = None) -> None:
    """多张图片合成一个 PDF（fpdf2，兼容 extract-video-ppt 的用法）。"""
    from fpdf import FPDF

    if not image_paths:
        raise ValueError("无图片可合成")
    pdf = FPDF()
    pdf.compress = False
    # 默认横向、按第一张图尺寸；或传入 size_wh (w, h)
    for img_path in image_paths:
        # 简单按单页一图；fpdf2 旧 API add_page(format=(h,w)) 横向
        if size_wh:
            w, h = size_wh
            pdf.add_page(orientation="L", format=(h, w))
        else:
            pdf.add_page(orientation="L")
        pdf.image(name=str(img_path), x=0, y=0, w=pdf.w, h=pdf.h)
    # 旧 API: output(path, "F")
    pdf.output(str(pdf_path), "F")


def frames_to_pptx(
    image_paths: list[Path],
    pptx_path: Path,
    *,
    slide_notes: list[str] | None = None,
) -> None:
    """多张图片生成 PowerPoint：一帧一页，可编辑。slide_notes 为每页备注（可选，预留语音转文字）。"""
    from pptx import Presentation
    from pptx.util import Inches

    if not image_paths:
        raise ValueError("无图片可合成")
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]  # Blank
    notes_list = slide_notes if slide_notes and len(slide_notes) >= len(image_paths) else None
    for i, img_path in enumerate(image_paths):
        slide = prs.slides.add_slide(blank)
        slide.shapes.add_picture(
            str(img_path),
            Inches(0), Inches(0),
            width=prs.slide_width,
            height=prs.slide_height,
        )
        if notes_list and notes_list[i].strip():
            slide.notes_slide.notes_text_frame.text = notes_list[i].strip()
    prs.save(str(pptx_path))
